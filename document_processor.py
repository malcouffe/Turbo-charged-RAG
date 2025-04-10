from typing import List, Optional, Dict, Any, Tuple
import os
import uuid
from datetime import datetime

from langchain.docstore.document import Document as LangchainDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings

# Pour les PDF
from langchain_community.document_loaders import PyPDFLoader
# Pour les PPTX
from pptx import Presentation
from docx import Document
import os
import base64
from datetime import datetime
from openai import OpenAI
from PIL import Image
import io
from image_analysis_cache import ImageAnalysisCache

class DocumentVectorStorePipeline:
    def __init__(
        self,
        openai_api_key: str,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        persist_directory: Optional[str] = None,
        image_output_dir: Optional[str] = "extracted_images",
        use_cache: bool = True
    ):
        """
        Initialise la pipeline pour PDF, PPTX et DOCX uniquement.

        Args:
            openai_api_key: API key for OpenAI embeddings
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
            persist_directory: Directory to persist vector store (optional)
            image_output_dir: Directory to save extracted images
            use_cache: Whether to use image analysis cache
        """
        self.embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            is_separator_regex=False
        )
        self.persist_directory = persist_directory
        self.image_output_dir = image_output_dir
        self.use_cache = use_cache
        self.image_cache = ImageAnalysisCache() if use_cache else None
        os.makedirs(image_output_dir, exist_ok=True)

    def analyze_image(self, image_path: str) -> str:
        """
        Use GPT-4 Vision to analyze an image and return a textual description.
        
        Args:
            image_path: Path to the image file
            
        Returns:
            str: Detailed description of the image
        """
        # Check cache first if enabled
        if self.use_cache and self.image_cache:
            cached_analysis = self.image_cache.get_analysis(image_path)
            if cached_analysis:
                return cached_analysis
        
        # Read and encode image
        with Image.open(image_path) as img:
            # Resize image if it's too large (max 2048px on longest side)
            if max(img.size) > 2048:
                ratio = 2048 / max(img.size)
                new_size = tuple(int(dim * ratio) for dim in img.size)
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to RGB if necessary
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Save to bytes
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG')
            img_byte_arr = img_byte_arr.getvalue()
            base64_image = base64.b64encode(img_byte_arr).decode('utf-8')

        # Get description from GPT-4 Vision
        response = self.openai_client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Please provide a detailed description of this image, including any visible text, key elements, and their relationships."},
                        {
                            "type": "image_url",
                            "image_url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    ]
                }
            ],
            max_tokens=300
        )
        
        analysis = response.choices[0].message.content
        
        # Save to cache if enabled
        if self.use_cache and self.image_cache:
            self.image_cache.add_analysis(image_path, analysis)
        
        return analysis

    def process_image_for_vectorstore(self, image_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Process an image and add its analysis to metadata.
        
        Args:
            image_metadata: Dictionary containing image information
            
        Returns:
            Dict with updated metadata including image analysis
        """
        image_path = image_metadata["image_path"]
        image_analysis = self.analyze_image(image_path)
        
        return {
            **image_metadata,
            "analysis": image_analysis
        }

    def extract_images(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extract images from PDF using markitdown.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of dictionaries containing image data and metadata
        """
        processor = PdfProcessor()
        images = []
        
        # Process the PDF and extract images
        pdf_content = processor.process_pdf(pdf_path)
        
        for idx, image in enumerate(pdf_content.images):
            # Generate unique filename for the image
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            image_filename = f"image_{timestamp}_{idx}.png"
            image_path = os.path.join(self.image_output_dir, image_filename)
            
            # Save the image
            with open(image_path, "wb") as f:
                f.write(base64.b64decode(image.base64_data))
            
            # Create image metadata
            image_metadata = {
                "image_path": image_path,
                "page_number": image.page_number,
                "width": image.width,
                "height": image.height,
                "caption": image.caption if hasattr(image, 'caption') else None
            }
            images.append(image_metadata)
        
        return images

    def process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Traite un fichier PPTX."""
        prs = Presentation(file_path)
        chunks = []
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            if slide_text.strip():
                chunks.append({
                    "content": slide_text.strip(),
                    "metadata": {"source": file_path, "slide_number": slide_number}
                })
        return chunks

    def process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Traite un fichier DOCX."""
        doc = docx.Document(file_path)
        chunks = []
        for para_number, paragraph in enumerate(doc.paragraphs, 1):
            if paragraph.text.strip():
                chunks.append({
                    "content": paragraph.text.strip(),
                    "metadata": {"source": file_path, "paragraph_number": para_number}
                })
        return chunks

    def process_document(self, file_path: str) -> Tuple[Chroma, List[LangchainDocument]]:
        """
        Traite un document (PDF, PPTX ou DOCX), crée des objets Document,
        les découpe en morceaux puis construit le vector store.
        
        Returns:
            Un tuple contenant le vector store et la liste des objets Document.
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        if file_extension == '.pdf':
            chunks = self.process_pdf(file_path)
        elif file_extension == '.pptx':
            chunks = self.process_pptx(file_path)
        elif file_extension == '.docx':
            chunks = self.process_docx(file_path)
        else:
            raise ValueError(f"Type de fichier non supporté: {file_extension}")

        # Crée des objets Document LangChain à partir des chunks
        documents = [
            LangchainDocument(page_content=chunk["content"], metadata=chunk["metadata"])
            for chunk in chunks
        ]
        # Découpage supplémentaire si nécessaire
        final_docs = self.text_splitter.split_documents(documents)
        # Assigner un id à chaque document si non défini
        for doc in final_docs:
            if not hasattr(doc, "id") or doc.id is None:
                doc.id = str(uuid.uuid4())
        # Construire le vector store via Chroma
        vector_store = Chroma.from_documents(
            documents=final_docs,
            embedding=self.embeddings,
            persist_directory=self.persist_directory
        )
        if self.persist_directory:
            vector_store.persist()  # Attention: méthode dépréciée depuis Chroma 0.4.x (les docs sont persistés automatiquement)
        return vector_store, final_docs

    def process_directory(self, directory_path: str) -> Tuple[Chroma, Dict[str, List[LangchainDocument]]]:
        """
        Traite tous les documents supportés dans un répertoire et combine les documents.

        Returns:
            Un tuple contenant le vector store combiné et un dictionnaire
            associant le chemin de chaque fichier à sa liste de documents.
        """
        all_docs: List[LangchainDocument] = []
        docs_by_file = {}
        supported_extensions = {'.pdf', '.pptx', '.docx'}

        for filename in os.listdir(directory_path):
            file_extension = os.path.splitext(filename)[1].lower()
            if file_extension in supported_extensions:
                file_path = os.path.join(directory_path, filename)
                vector_store, docs = self.process_document(file_path)
                docs_by_file[file_path] = docs
                all_docs.extend(docs)

        combined_vector_store = Chroma.from_documents(
            documents=all_docs,
            embedding=self.embeddings,
            persist_directory=self.persist_directory
        )
        if self.persist_directory:
            vector_store.persist()
            
        return vector_store, all_images

    def merge_vector_stores(self, existing_store: Chroma, new_store: Chroma) -> Chroma:
        """
        Merge a new vector store into an existing one.
        
        Args:
            existing_store: The existing Chroma vector store
            new_store: The new Chroma vector store to merge
            
        Returns:
            The merged Chroma vector store
        """
        # Get documents from the new vector store
        new_docs = []
        for doc_id in new_store._collection.get()["ids"]:
            result = new_store._collection.get([doc_id])
            new_docs.append({
                "document": result["documents"][0],
                "metadata": result["metadatas"][0],
                "embedding": result["embeddings"][0]
            })
        
        # Add new documents to existing store
        for doc in new_docs:
            existing_store._collection.add(
                documents=[doc["document"]],
                embeddings=[doc["embedding"]],
                metadatas=[doc["metadata"]]
            )
        
        # Persist if needed
        if self.persist_directory:
            existing_store.persist()
            
        return existing_store